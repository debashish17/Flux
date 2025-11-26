from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import re

def parse_markdown_line(paragraph, text):
    """
    Parse a line of markdown text and add formatted runs to a paragraph
    Handles: **bold**, *italic*, `code`, [links](url)
    """
    from docx.shared import RGBColor

    # Pattern to match markdown formatting
    # Matches: **bold**, *italic*, `code`, [text](url)
    pattern = r'(\*\*.*?\*\*|\*.*?\*|`.*?`|\[.*?\]\(.*?\))'

    parts = re.split(pattern, text)

    for part in parts:
        if not part:
            continue

        # Bold: **text**
        if part.startswith('**') and part.endswith('**'):
            text_content = part[2:-2]
            run = paragraph.add_run(text_content)
            run.bold = True

        # Italic: *text*
        elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
            text_content = part[1:-1]
            run = paragraph.add_run(text_content)
            run.italic = True

        # Code: `text`
        elif part.startswith('`') and part.endswith('`'):
            text_content = part[1:-1]
            run = paragraph.add_run(text_content)
            run.font.name = 'Courier New'
            run.font.color.rgb = RGBColor(220, 38, 38)  # Red color for code

        # Link: [text](url)
        elif part.startswith('[') and '](' in part:
            match = re.match(r'\[(.*?)\]\((.*?)\)', part)
            if match:
                link_text = match.group(1)
                run = paragraph.add_run(link_text)
                run.font.color.rgb = RGBColor(99, 102, 241)  # Blue color for links
                run.underline = True
        else:
            # Regular text
            paragraph.add_run(part)


def create_docx(project) -> io.BytesIO:
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from datetime import datetime

    doc = Document()

    # Set page margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.25)
        section.right_margin = Inches(1.25)

    # Set default font for the document
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # ============================================
    # TITLE PAGE
    # ============================================
    # Add spacing from top
    for _ in range(8):
        doc.add_paragraph()

    # Add title - centered, large, professional
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(project.title)
    title_run.font.name = 'Calibri'
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(31, 78, 120)

    # Add spacing
    doc.add_paragraph()
    doc.add_paragraph()

    # Add decorative line
    line_para = doc.add_paragraph()
    line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    line_run = line_para.add_run('─' * 40)
    line_run.font.color.rgb = RGBColor(100, 100, 100)

    # Add spacing
    doc.add_paragraph()
    doc.add_paragraph()

    # Add subtitle/date
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_run = date_para.add_run(f"Generated on {datetime.now().strftime('%B %d, %Y')}")
    date_run.font.name = 'Calibri'
    date_run.font.size = Pt(12)
    date_run.font.color.rgb = RGBColor(100, 100, 100)

    # Page break after title page
    doc.add_page_break()

    # ============================================
    # TABLE OF CONTENTS PAGE
    # ============================================
    toc_heading = doc.add_heading('Table of Contents', level=1)
    toc_heading_run = toc_heading.runs[0]
    toc_heading_run.font.name = 'Calibri'
    toc_heading_run.font.size = Pt(20)
    toc_heading_run.font.color.rgb = RGBColor(31, 78, 120)

    doc.add_paragraph()

    # Add TOC entries (extract from sections)
    toc_entries = []
    for section in sorted(project.sections, key=lambda x: x.orderIndex):
        content = section.content or ""
        # Extract main headings from markdown
        if '##' in content:
            lines = content.split('\n')
            for line in lines:
                stripped = line.strip()
                if stripped.startswith('## ') and not stripped.startswith('###'):
                    heading_text = stripped[3:].strip()
                    toc_entries.append((heading_text, 1))
                elif stripped.startswith('### '):
                    heading_text = stripped[4:].strip()
                    toc_entries.append((heading_text, 2))

    # Generate TOC
    for entry_text, level in toc_entries:
        toc_para = doc.add_paragraph()
        if level == 1:
            toc_para.paragraph_format.left_indent = Inches(0.25)
            toc_run = toc_para.add_run(f"• {entry_text}")
            toc_run.font.size = Pt(12)
            toc_run.font.bold = True
        else:  # level 2
            toc_para.paragraph_format.left_indent = Inches(0.75)
            toc_run = toc_para.add_run(f"◦ {entry_text}")
            toc_run.font.size = Pt(11)
        toc_run.font.name = 'Calibri'
        toc_run.font.color.rgb = RGBColor(68, 114, 196)
        toc_para.paragraph_format.space_after = Pt(6)

    # Page break after TOC
    doc.add_page_break()

    for section in sorted(project.sections, key=lambda x: x.orderIndex):
        content = section.content or "No content generated yet."

        # Check if content is Markdown (contains ## headings or markdown formatting)
        is_markdown = '##' in content or '**' in content or '*' in content

        if is_markdown:
            # Process Markdown content
            lines = content.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].strip()

                if not line:
                    i += 1
                    continue

                # Heading level 2: ## Text
                if line.startswith('## '):
                    heading_text = line[3:].strip()
                    heading = doc.add_heading(heading_text, level=1)
                    heading_run = heading.runs[0] if heading.runs else None
                    if heading_run:
                        heading_run.font.name = 'Calibri'
                        heading_run.font.size = Pt(16)
                        heading_run.font.color.rgb = RGBColor(68, 114, 196)
                    heading.space_before = Pt(12)
                    heading.space_after = Pt(6)

                # Heading level 3: ### Text
                elif line.startswith('### '):
                    heading_text = line[4:].strip()
                    heading = doc.add_heading(heading_text, level=2)
                    heading_run = heading.runs[0] if heading.runs else None
                    if heading_run:
                        heading_run.font.name = 'Calibri'
                        heading_run.font.size = Pt(13)
                        heading_run.font.color.rgb = RGBColor(89, 89, 89)
                    heading.space_before = Pt(10)
                    heading.space_after = Pt(4)

                # Heading level 4: #### Text
                elif line.startswith('#### '):
                    heading_text = line[5:].strip()
                    heading = doc.add_heading(heading_text, level=3)
                    heading_run = heading.runs[0] if heading.runs else None
                    if heading_run:
                        heading_run.font.name = 'Calibri'
                        heading_run.font.size = Pt(12)
                        heading_run.font.color.rgb = RGBColor(112, 112, 112)
                    heading.space_before = Pt(8)
                    heading.space_after = Pt(3)

                # Unordered list: - item or * item
                elif line.startswith('- ') or line.startswith('* '):
                    bullet_text = line[2:].strip()
                    p = doc.add_paragraph(style='List Bullet')
                    parse_markdown_line(p, bullet_text)
                    p.paragraph_format.left_indent = Inches(0.25)
                    p.paragraph_format.space_after = Pt(6)

                # Ordered list: 1. item
                elif re.match(r'^\d+\.\s', line):
                    list_text = re.sub(r'^\d+\.\s', '', line).strip()
                    p = doc.add_paragraph(style='List Number')
                    parse_markdown_line(p, list_text)
                    p.paragraph_format.left_indent = Inches(0.25)
                    p.paragraph_format.space_after = Pt(6)

                # Blockquote: > text
                elif line.startswith('> '):
                    quote_text = line[2:].strip()
                    p = doc.add_paragraph()
                    parse_markdown_line(p, quote_text)
                    p.paragraph_format.left_indent = Inches(0.5)
                    p.paragraph_format.space_after = Pt(10)
                    # Add shading for blockquotes
                    from docx.oxml.shared import OxmlElement
                    from docx.oxml.ns import qn
                    shading_elm = OxmlElement('w:shd')
                    shading_elm.set(qn('w:fill'), 'F0F0F0')
                    p._element.get_or_add_pPr().append(shading_elm)

                # Table: | Column | Column |
                elif line.startswith('|') and '|' in line[1:]:
                    # Start of table - collect all table rows
                    table_rows = []
                    table_rows.append(line)
                    i += 1

                    # Collect subsequent table rows
                    while i < len(lines):
                        next_line = lines[i].strip()
                        if next_line.startswith('|'):
                            table_rows.append(next_line)
                            i += 1
                        else:
                            i -= 1  # Step back to process this line normally
                            break

                    # Parse and create table
                    if len(table_rows) >= 2:  # At least header + separator
                        # Parse header
                        header_cells = [cell.strip() for cell in table_rows[0].split('|')[1:-1]]

                        # Skip separator row (row with dashes)
                        data_rows = []
                        for row in table_rows[2:]:  # Skip header and separator
                            cells = [cell.strip() for cell in row.split('|')[1:-1]]
                            if cells and any(cells):  # Not empty row
                                data_rows.append(cells)

                        # Create table in document
                        if header_cells and data_rows:
                            table = doc.add_table(rows=1 + len(data_rows), cols=len(header_cells))
                            table.style = 'Light Grid Accent 1'

                            # Set header row
                            header_row = table.rows[0]
                            for idx, cell_text in enumerate(header_cells):
                                cell = header_row.cells[idx]
                                cell.text = cell_text
                                # Bold header text
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(11)

                            # Set data rows
                            for row_idx, row_data in enumerate(data_rows):
                                table_row = table.rows[row_idx + 1]
                                for col_idx, cell_text in enumerate(row_data):
                                    if col_idx < len(header_cells):
                                        table_row.cells[col_idx].text = cell_text

                            # Add spacing after table
                            doc.add_paragraph()

                # Regular paragraph
                else:
                    p = doc.add_paragraph()
                    parse_markdown_line(p, line)
                    p.paragraph_format.space_after = Pt(10)
                    p.paragraph_format.line_spacing = 1.15
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                i += 1
        else:
            # Legacy format: Plain text processing
            heading = doc.add_heading(section.title, level=1)
            heading_run = heading.runs[0]
            heading_run.font.name = 'Calibri'
            heading_run.font.size = Pt(16)
            heading_run.font.color.rgb = RGBColor(68, 114, 196)
            heading.space_before = Pt(12)
            heading.space_after = Pt(6)

            paragraphs = content.split('\n\n')
            for para_text in paragraphs:
                para_text = para_text.strip()
                if para_text:
                    if para_text.startswith(('•', '-', '*')):
                        lines = para_text.split('\n')
                        for line in lines:
                            line = line.strip()
                            if line:
                                bullet_text = line.lstrip('•-*').strip()
                                if bullet_text:
                                    p = doc.add_paragraph(bullet_text, style='List Bullet')
                                    p.paragraph_format.left_indent = Inches(0.25)
                                    p.paragraph_format.space_after = Pt(6)
                    else:
                        p = doc.add_paragraph(para_text)
                        p.paragraph_format.space_after = Pt(10)
                        p.paragraph_format.line_spacing = 1.15
                        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return file_stream

def parse_slide_content(content: str) -> dict:
    """Parse the structured slide content format"""
    result = {
        'title': '',
        'bullets': [],
        'image_suggestion': ''
    }

    if not content:
        return result

    lines = content.split('\n')
    in_content_section = False

    # Keywords that indicate speaker notes (to be filtered out)
    speaker_note_keywords = ['SPEAKER_NOTES:', 'NOTES:', 'SPEAKER:', 'NOTE:']

    for line in lines:
        line = line.strip()

        # Skip any lines that are speaker notes
        if any(line.upper().startswith(keyword) for keyword in speaker_note_keywords):
            in_content_section = False
            continue

        if line.startswith('TITLE:'):
            result['title'] = line.replace('TITLE:', '').strip()
        elif line.startswith('CONTENT:'):
            in_content_section = True
        elif line.startswith('IMAGE_SUGGESTION:'):
            result['image_suggestion'] = line.replace('IMAGE_SUGGESTION:', '').strip()
            in_content_section = False
        elif in_content_section and line:
            # Extract bullet points
            bullet = line.lstrip('•-*').strip()
            if bullet:
                result['bullets'].append(bullet)

    return result

def create_pptx(project) -> io.BytesIO:
    prs = Presentation()

    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Format title
    title.text = project.title
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(31, 78, 120)

    subtitle.text = "AI-Generated Presentation"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(100, 100, 100)

    # Content Slides
    for section in sorted(project.sections, key=lambda x: x.orderIndex):
        # Use blank layout for more control
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Parse the structured content
        parsed = parse_slide_content(section.content or "")

        # Add title - matching the editor's style (text-4xl font-bold text-gray-900)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = parsed['title'] or section.title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)  # Matches text-4xl
        title_para.font.bold = True
        title_para.font.name = 'Calibri'
        title_para.font.color.rgb = RGBColor(17, 24, 39)  # text-gray-900

        # Determine layout based on whether there's an image suggestion
        if parsed['image_suggestion']:
            # Two-column layout: content on left, image placeholder on right
            content_left = Inches(0.5)
            content_top = Inches(1.5)
            content_width = Inches(5)
            content_height = Inches(3.5)

            # Add content textbox
            content_box = slide.shapes.add_textbox(
                content_left, content_top, content_width, content_height
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Add bullet points - matching editor style with actual bullets
            for i, bullet in enumerate(parsed['bullets']):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                # Set bullet text
                p.text = bullet
                p.level = 0

                # Add bullet using XML directly
                from pptx.oxml.xmlchemy import OxmlElement
                pPr = p._element.get_or_add_pPr()

                # Create bullet node
                buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                if buNone is not None:
                    pPr.remove(buNone)

                # Add bullet character
                buChar = OxmlElement('a:buChar')
                buChar.set('char', '•')
                pPr.append(buChar)

                # Style the text
                p.font.size = Pt(18)  # text-lg
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(55, 65, 81)  # text-gray-700
                p.space_before = Pt(12)  # space-y-4

            # Add image placeholder
            img_left = Inches(6)
            img_top = Inches(1.5)
            img_width = Inches(3.5)
            img_height = Inches(3.5)

            img_placeholder = slide.shapes.add_textbox(
                img_left, img_top, img_width, img_height
            )
            img_frame = img_placeholder.text_frame
            img_frame.word_wrap = True  # Enable word wrapping
            img_frame.vertical_anchor = 1  # Middle vertical alignment

            # Add image icon and text
            img_frame.text = f"📷 Image:\n{parsed['image_suggestion']}"
            img_para = img_frame.paragraphs[0]
            img_para.font.size = Pt(12)
            img_para.font.italic = True
            img_para.font.color.rgb = RGBColor(107, 114, 128)  # Gray
            img_para.alignment = 1  # Center

            # Add border to image placeholder
            line = img_placeholder.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Pt(1)
        else:
            # Full-width content layout
            content_left = Inches(1)
            content_top = Inches(1.5)
            content_width = Inches(8)
            content_height = Inches(3.5)

            content_box = slide.shapes.add_textbox(
                content_left, content_top, content_width, content_height
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Add bullet points - matching editor style with actual bullets
            for i, bullet in enumerate(parsed['bullets']):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                # Set bullet text
                p.text = bullet
                p.level = 0

                # Add bullet using XML directly
                from pptx.oxml.xmlchemy import OxmlElement
                pPr = p._element.get_or_add_pPr()

                # Create bullet node
                buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                if buNone is not None:
                    pPr.remove(buNone)

                # Add bullet character
                buChar = OxmlElement('a:buChar')
                buChar.set('char', '•')
                pPr.append(buChar)

                # Style the text
                p.font.size = Pt(20)
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(55, 65, 81)  # text-gray-700
                p.space_before = Pt(12)  # space-y-4

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream
